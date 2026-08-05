"""
Chargeur de fichiers PPTX pour le moteur d'intelligence documentaire GalSen IA.
"""

import os
from .document_loader_factory import BaseDocumentLoader
from .types import DocumentItem, DocumentType
import logging

logger = logging.getLogger(__name__)


class PPTXLoader(BaseDocumentLoader):
    """Chargeur pour les fichiers PPTX (.pptx)."""

    # Extensions supportées
    supported_extensions = ['.pptx']

    def load(self, source: str, **kwargs) -> DocumentItem:
        """
        Charge un fichier PPTX et retourne un DocumentItem.

        Args:
            source: Chemin vers le fichier PPTX
            **kwargs: Options supplémentaires

        Returns:
            DocumentItem représentant le fichier chargé
        """
        if not os.path.isfile(source):
            raise FileNotFoundError(f"Fichier non trouvé: {source}")

        try:
            import pptx
        except ImportError:
            message = (
                "python-pptx est requis pour charger les fichiers PPTX. "
                "Installez-le avec: pip install python-pptx"
            )
            logger.error(message)
            raise ImportError(message)

        # Charger la présentation
        presentation = pptx.Presentation(source)

        # Extraire le texte de toutes les diapositives
        text_content = []
        for slide_num, slide in enumerate(presentation.slides):
            text_content.append(f"Diapositive {slide_num + 1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            text_content.append("")  # Ligne vide entre les diapositives

        full_text = '\n'.join(text_content)

        # Extraire les métadonnées de la présentation
        props = presentation.core_properties
        metadata = {
            "size": os.path.getsize(source),
            "created": os.path.getctime(source),
            "modified": os.path.getmtime(source),
            "source": os.path.abspath(source),
            "filename": os.path.basename(source),
            "extension": ".pptx",
            "slide_count": len(presentation.slides)
        }

        if props.author:
            metadata["author"] = props.author
        if props.created:
            metadata["created"] = props.created.timestamp()
        if props.modified:
            metadata["modified"] = props.modified.timestamp()
        if props.title:
            metadata["title"] = props.title
        if props.subject:
            metadata["subject"] = props.subject
        if hasattr(props, 'keywords') and props.keywords:
            # Les mots-clés peuvent être une chaîne séparée par des virgules
            if isinstance(props.keywords, str):
                metadata["keywords"] = [k.strip() for k in props.keywords.split(',')]
            else:
                metadata["keywords"] = list(props.keywords)
        if props.category:
            metadata["category"] = props.category

        # Utiliser le titre de la présentation ou le nom du fichier
        title = metadata.get("title") or os.path.splitext(os.path.basename(source))[0]

        # Créer l'élément de document
        document = DocumentItem(
            document_id="",
            document_type=DocumentType.PPTX,
            title=title,
            content=full_text,
            metadata=metadata
        )

        return document

    def supports_type(self, document_type: str) -> bool:
        """Vérifie si le chargeur supporte le type de document donné."""
        return document_type in [
            DocumentType.PPTX.value,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        ]